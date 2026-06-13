import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Scheme
COLOR_PRIMARY_NAVY = RGBColor(13, 43, 69)
COLOR_ACCENT_TEAL = RGBColor(32, 178, 170)
COLOR_MUTED_GRAY = RGBColor(100, 110, 120)
COLOR_BG_LIGHT = RGBColor(248, 249, 250)
COLOR_WHITE = RGBColor(255, 255, 255)

def apply_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, category_text="PEARLYPATH AI"):
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = "Arial"
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_ACCENT_TEAL
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Arial"
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_PRIMARY_NAVY

def try_insert_logo(slide, left, top, width):
    target_path = os.path.join("PearlyPath AI", "public", "logo.jpg")
    if os.path.exists(target_path):
        slide.shapes.add_picture(target_path, left, top, width=width)
    elif os.path.exists("logo.jpg"):
        slide.shapes.add_picture("logo.jpg", left, top, width=width)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_MUTED_GRAY
        tx_frame = shape.text_frame
        tx_frame.text = "[ PearlyPath Logo Place ]"

# SLIDE 1
blank_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(blank_layout)
apply_background(slide1, COLOR_PRIMARY_NAVY)

title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.5))
tf1 = title_box.text_frame
tf1.word_wrap = True
p1 = tf1.paragraphs[0]
p1.text = "PEARLYPATH AI"
p1.font.name = "Arial"
p1.font.size = Pt(54)
p1.font.bold = True
p1.font.color.rgb = COLOR_WHITE

p2 = tf1.add_paragraph()
p2.text = "Bahrain’s MOST Intelligent Tourism Companion"
p2.font.name = "Arial"
p2.font.size = Pt(22)
p2.font.color.rgb = COLOR_ACCENT_TEAL
p2.space_before = Pt(14)

p3 = tf1.add_paragraph()
p3.text = '"Your AI Guide to Bahrain"'
p3.font.name = "Arial"
p3.font.size = Pt(18)
p3.font.italic = True
p3.font.color.rgb = COLOR_WHITE
p3.space_before = Pt(8)

meta_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(7.0), Inches(1.5))
tf_meta = meta_box.text_frame
tf_meta.word_wrap = True
p_team = tf_meta.paragraphs[0]
p_team.text = "Team: Afroos, Dev, Faaiz, Rehan, Sharan"
p_team.font.size = Pt(14)
p_team.font.color.rgb = COLOR_WHITE

p_venue = tf_meta.add_paragraph()
p_venue.text = "STREAMSCAPES 2026"
p_venue.font.size = Pt(14)
p_venue.font.bold = True
p_venue.font.color.rgb = COLOR_ACCENT_TEAL
p_venue.space_before = Pt(6)

try_insert_logo(slide1, Inches(9.5), Inches(4.5), Inches(2.8))

# SLIDE 2
slide2 = prs.slides.add_slide(blank_layout)
apply_background(slide2, COLOR_BG_LIGHT)
add_header(slide2, "Tourism Challenges in Bahrain", "The Problem")

left_col = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.5))
tf_left = left_col.text_frame
tf_left.word_wrap = True
p = tf_left.paragraphs[0]
p.text = "Tourists Often Face:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY_NAVY

for ch in ["Struggle to find personalized experiences", "Visit only hyper-popular destinations", "Miss out on authentic cultural attractions", "Face friction with local language & planning barriers"]:
    p = tf_left.add_paragraph()
    p.text = f"❌  {ch}"
    p.font.size = Pt(14)
    p.space_before = Pt(12)
    p.font.color.rgb = COLOR_PRIMARY_NAVY

right_col = slide2.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.5))
tf_right = right_col.text_frame
tf_right.word_wrap = True
p = tf_right.paragraphs[0]
p.text = "Local Businesses Often:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY_NAVY

for bz in ["Small local businesses receive limited tourist visibility", "Tourism spending remains concentrated in a few locations"]:
    p = tf_right.add_paragraph()
    p.text = f"⚠️  {bz}"
    p.font.size = Pt(14)
    p.space_before = Pt(12)
    p.font.color.rgb = COLOR_PRIMARY_NAVY

showcase = slide2.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.5))
tf_show = showcase.text_frame
tf_show.word_wrap = True
p = tf_show.paragraphs[0]
p.text = "Popular Mainstream Tracks: The Avenues • Bahrain City Centre • Amwaj Islands"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = COLOR_MUTED_GRAY
p = tf_show.add_paragraph()
p.text = "Overlooked Heritage Gems: Rashid's Cafeteria • Al Jasra Handicrafts Centre • Traditional Halwa Shops • Old Textile Souq"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT_TEAL
p.space_before = Pt(4)

# SLIDE 3
slide3 = prs.slides.add_slide(blank_layout)
apply_background(slide3, COLOR_BG_LIGHT)
add_header(slide3, "Introducing PearlyPath AI", "Our Solution")

intro_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.0))
tf_intro = intro_box.text_frame
tf_intro.word_wrap = True
p = tf_intro.paragraphs[0]
p.text = "An AI-Powered Tourism Companion designed to engineer personalized travel experiences while actively expanding visibility for Bahrain's cultural heritage markers and native enterprises."
p.font.size = Pt(16)
p.font.color.rgb = COLOR_PRIMARY_NAVY

in_shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.8), Inches(3.6), Inches(3.4))
in_shape.fill.solid()
in_shape.fill.fore_color.rgb = COLOR_WHITE
in_shape.line.color.rgb = COLOR_MUTED_GRAY
tf_in = in_shape.text_frame
tf_in.word_wrap = True
p = tf_in.paragraphs[0]
p.text = "USER INPUT ENGINE"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT_TEAL
for inp in ["• User Mood Data", "• Strict Budget Inputs", "• Trip Duration Bounds", "• Granular Interests"]:
    p = tf_in.add_paragraph()
    p.text = inp
    p.font.size = Pt(14)
    p.space_before = Pt(14)
    p.font.color.rgb = COLOR_PRIMARY_NAVY

core_shape = slide3.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(4.8), Inches(4.0), Inches(3.6), Inches(1.0))
core_shape.fill.solid()
core_shape.fill.fore_color.rgb = COLOR_PRIMARY_NAVY
tf_core = core_shape.text_frame
p = tf_core.paragraphs[0]
p.text = "PEARLYPATH ENGINE"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = COLOR_WHITE
p.alignment = PP_ALIGN.CENTER

out_shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(2.8), Inches(3.7), Inches(3.4))
out_shape.fill.solid()
out_shape.fill.fore_color.rgb = COLOR_WHITE
out_shape.line.color.rgb = COLOR_ACCENT_TEAL
tf_out = out_shape.text_frame
tf_out.word_wrap = True
p = tf_out.paragraphs[0]
p.text = "DYNAMIC OUTPUTS"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT_TEAL
for outp in ["• Personalized Itineraries", "• Hidden Gem Matrix", "• Cultural Guidelines", "• Real-time Translation"]:
    p = tf_out.add_paragraph()
    p.text = outp
    p.font.size = Pt(14)
    p.space_before = Pt(14)
    p.font.color.rgb = COLOR_PRIMARY_NAVY

# SLIDE 4
slide4 = prs.slides.add_slide(blank_layout)
apply_background(slide4, COLOR_BG_LIGHT)
add_header(slide4, "Smart Tourism Core Architecture", "Key Features")

features_data = [
    {"title": "AI Itinerary Planner", "desc": "Compiles strict multi-point schedules handling custom budgets, tailored personal interests, and duration lengths seamlessly."},
    {"title": "AI Tourism Assistant", "desc": "Instantly processes complex native cultural etiquette queries, transit schedules, and deep-layer restaurant lookups."},
    {"title": "Crowd Management", "desc": "Leverages predictive modeling to forecast destination congestion, preserving comfortable visitor volume layouts."},
    {"title": "Hidden Gem Router", "desc": "Directs traffic safely to Monet Cafe, Al Jasra, and local Souqs instead of only over-indexing mainstream hubs."}
]

for idx in range(len(features_data)):
    feat = features_data[idx]
    current_x = Inches(0.8) + idx * (Inches(2.7) + Inches(0.3))
    panel = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, current_x, Inches(2.2), Inches(2.7), Inches(4.4))
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLOR_WHITE
    panel.line.color.rgb = COLOR_MUTED_GRAY
    tf_p = panel.text_frame
    tf_p.word_wrap = True
    p = tf_p.paragraphs[0]
    p.text = feat["title"]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_NAVY
    p_desc = tf_p.add_paragraph()
    p_desc.text = feat["desc"]
    p_desc.font.size = Pt(13)
    p_desc.font.color.rgb = COLOR_MUTED_GRAY
    p_desc.space_before = Pt(14)

# SLIDE 5
slide5 = prs.slides.add_slide(blank_layout)
apply_background(slide5, COLOR_BG_LIGHT)
add_header(slide5, "Operational Execution Pipeline", "How It Works")

steps = [
    {"num": "01", "name": "Input Capture", "details": "User profiles mood, budgetary parameters, and specific cultural interests."},
    {"num": "02", "name": "Context Analysis", "details": "The AI core contextualizes native geographic nodes and proximity constraints."},
    {"num": "03", "name": "Engine Selection", "details": "Algorithm matches optimized paths connecting overlooked local vendors."},
    {"num": "04", "name": "Live Delivery", "details": "User renders live dynamic routes, translation bridges, and custom maps."}
]

for idx in range(len(steps)):
    st = steps[idx]
    cx = Inches(0.8) + idx * (Inches(2.7) + Inches(0.3))
    box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(2.5), Inches(2.7), Inches(3.8))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_PRIMARY_NAVY if idx == 0 else COLOR_WHITE
    box.line.color.rgb = COLOR_ACCENT_TEAL
    tf_b = box.text_frame
    tf_b.word_wrap = True
    p = tf_b.paragraphs[0]
    p.text = st["num"]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_TEAL if idx != 0 else COLOR_WHITE
    p_name = tf_b.add_paragraph()
    p_name.text = st["name"]
    p_name.font.size = Pt(16)
    p_name.font.bold = True
    p_name.font.color.rgb = COLOR_PRIMARY_NAVY if idx != 0 else COLOR_WHITE
    p_name.space_before = Pt(8)
    p_det = tf_b.add_paragraph()
    p_det.text = st["details"]
    p_det.font.size = Pt(12)
    p_det.font.color.rgb = COLOR_MUTED_GRAY if idx != 0 else COLOR_WHITE
    p_det.space_before = Pt(10)

# SLIDE 6
slide6 = prs.slides.add_slide(blank_layout)
apply_background(slide6, COLOR_BG_LIGHT)
add_header(slide6, "Sustainable Architecture Revenue Streams", "Business Model")

models = [
    {"title": "B2B Vendor Partnerships", "desc": "Integrated native dining, hospitality venues, and specialty boutiques access high-intent premium placement inside recommendations loops dynamically."},
    {"title": "B2C Premium Membership", "desc": "Unlocks high-tier continuous mapping structures, completely offline local translation models, and advanced concierge optimization loops."},
    {"title": "Macro Tourism Analytics", "desc": "Aggregates absolute anonymous structural trend tracking for government agencies, hospitality entities, and systemic planning councils."}
]

for idx in range(len(models)):
    md = models[idx]
    mx = Inches(0.8) + idx * (Inches(3.7) + Inches(0.3))
    m_box = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, mx, Inches(2.4), Inches(3.7), Inches(4.0))
    m_box.fill.solid()
    m_box.fill.fore_color.rgb = COLOR_WHITE
    m_box.line.color.rgb = COLOR_ACCENT_TEAL if idx == 0 else COLOR_MUTED_GRAY
    tf_m = m_box.text_frame
    tf_m.word_wrap = True
    p = tf_m.paragraphs[0]
    p.text = md["title"]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_NAVY
    p_d = tf_m.add_paragraph()
    p_d.text = md["desc"]
    p_d.font.size = Pt(13)
    p_d.font.color.rgb = COLOR_MUTED_GRAY
    p_d.space_before = Pt(14)

# SLIDE 7
slide7 = prs.slides.add_slide(blank_layout)
apply_background(slide7, COLOR_BG_LIGHT)
add_header(slide7, "Impact on Bahrain's Tourism Ecosystem", "Benefits for Bahrain")

benefits = [
    {"title": "Empowering Local Micro-SMEs", "desc": "Drives immediate monetization directly to neighborhood cafes, independent handicraft artisans, and heritage sites."},
    {"title": "Elevating the Visitor UX Matrix", "desc": "Replaces generic travel stress with deeply contextualized native schedules and real-time navigation pipelines."},
    {"title": "Geographic De-congestion Mechanics", "desc": "Redirects volume patterns safely outward from the standard shopping strip hubs into authentic cultural destinations."},
    {"title": "Amplified Economic Circulation", "desc": "Longer exploration pathways directly correlate with diversified regional spending velocities across sectors."}
]

for idx in range(len(benefits)):
    bn = benefits[idx]
    bx = Inches(0.8) if idx % 2 == 0 else Inches(6.8)
    by = Inches(2.2) if idx < 2 else Inches(4.4)
    b_card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, Inches(5.6), Inches(1.8))
    b_card.fill.solid()
    b_card.fill.fore_color.rgb = COLOR_WHITE
    b_card.line.color.rgb = COLOR_MUTED_GRAY
    tf_b = b_card.text_frame
    tf_b.word_wrap = True
    p = tf_b.paragraphs[0]
    p.text = bn["title"]
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_NAVY
    p_ds = tf_b.add_paragraph()
    p_ds.text = bn["desc"]
    p_ds.font.size = Pt(12)
    p_ds.font.color.rgb = COLOR_MUTED_GRAY
    p_ds.space_before = Pt(6)

# SLIDE 8
slide8 = prs.slides.add_slide(blank_layout)
apply_background(slide8, COLOR_BG_LIGHT)
add_header(slide8, "Supporting SDG 8: Decent Work & Economic Growth", "UN SDG Alignment")

sdg_box = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.5))
tf_sdg = sdg_box.text_frame
tf_sdg.word_wrap = True

sdg_points = [
    ("Direct SME Capitalization:", "Establishes long-tail discoverability pipelines for highly vulnerable heritage artisan trades."),
    ("Tourism Volume Retention:", "Safely extends visitor retention timetables, directly elevating regional hospitality metrics."),
    ("Frontline Digital Innovation:", "Deploys production-grade AI algorithms to fundamentally upgrade and modernize the nation's service economy frameworks."),
    ("Sustainable Economic Scaling:", "Engineers scalable, high-yield infrastructure growth completely independent of heavy capital asset footprints.")
]

for title, desc in sdg_points:
    p = tf_sdg.add_paragraph() if tf_sdg.text else tf_sdg.paragraphs[0]
    p.text = f"✔  {title} "
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY_NAVY
    run = p.add_run()
    run.text = desc
    run.font.bold = False
    run.font.color.rgb = COLOR_MUTED_GRAY
    p.space_before = Pt(14)

callout = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.6), Inches(11.733), Inches(1.0))
callout.fill.solid()
callout.fill.fore_color.rgb = COLOR_PRIMARY_NAVY
callout.line.fill.background()
tf_call = callout.text_frame
p_call = tf_call.paragraphs[0]
p_call.text = "Key Engine Strategy: PearlyPath AI converts native tourism paths into long-term decentralized economic opportunities."
p_call.font.size = Pt(14)
p_call.font.bold = True
p_call.font.color.rgb = COLOR_WHITE
p_call.alignment = PP_ALIGN.CENTER

# SLIDE 9
slide9 = prs.slides.add_slide(blank_layout)
apply_background(slide9, COLOR_PRIMARY_NAVY)

info_box = slide9.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(6.5), Inches(4.5))
tf_info = info_box.text_frame
tf_info.word_wrap = True
p = tf_info.paragraphs[0]
p.text = "THANK YOU"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = COLOR_WHITE

p = tf_info.add_paragraph()
p.text = "PearlyPath AI"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT_TEAL
p.space_before = Pt(8)

p = tf_info.add_paragraph()
p.text = '"Discover Bahrain Intelligently."'
p.font.size = Pt(18)
p.font.italic = True
p.font.color.rgb = COLOR_WHITE
p.space_before = Pt(6)

p = tf_info.add_paragraph()
p.text = "Questions & Answers Session"
p.font.size = Pt(16)
p.font.color.rgb = COLOR_MUTED_GRAY
p.space_before = Pt(28)

qr_zone = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.8), Inches(3.8), Inches(4.2))
qr_zone.fill.solid()
qr_zone.fill.fore_color.rgb = COLOR_WHITE
qr_zone.line.fill.background()
tf_qr = qr_zone.text_frame
tf_qr.word_wrap = True
p = tf_qr.paragraphs[0]
p.text = "SCAN FOR LIVE DEMO"
p.font.size = Pt(13)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY_NAVY
p.alignment = PP_ALIGN.CENTER

p = tf_qr.add_paragraph()
p.text = "[ Drop QR Code Here ]"
p.font.size = Pt(12)
p.font.color.rgb = COLOR_MUTED_GRAY
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(44)

p = tf_qr.add_paragraph()
p.text = "📱 Scan the QR code to explore PearlyPath AI live on your mobile device."
p.font.size = Pt(11)
p.font.color.rgb = COLOR_PRIMARY_NAVY
p.alignment = PP_ALIGN.CENTER
p.space_before = Pt(54)

prs.save("PearlyPath_AI_Presentation.pptx")
print("Success: 'PearlyPath_AI_Presentation.pptx' created perfectly!")
